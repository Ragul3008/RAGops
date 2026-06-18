from pathlib import Path
import tempfile
import logging

logger = logging.getLogger(__name__)


class DocumentParser:
    """Multi-format document parser.

    Uses format-specific fast paths first (pypdf for PDF,
    python-docx for DOCX, plain decode for text) and falls
    back to unstructured.partition only when needed.

    Supported formats: .pdf  .docx  .pptx  .txt  .md  .html
    """

    SUPPORTED = {".pdf", ".docx", ".txt", ".md", ".html", ".pptx"}

    # ------------------------------------------------------------------ #
    #  Public API                                                          #
    # ------------------------------------------------------------------ #

    async def parse(self, content: bytes, filename: str) -> str:
        ext = Path(filename).suffix.lower()
        if ext not in self.SUPPORTED:
            raise ValueError(
                f"Unsupported format '{ext}'. "
                f"Supported: {', '.join(sorted(self.SUPPORTED))}"
            )

        try:
            return await self._dispatch(content, ext, filename)
        except ImportError as exc:
            # Surface a friendly message instead of a raw 500
            raise RuntimeError(
                f"A required dependency for '{ext}' files is not installed. "
                f"Run: pip install \"unstructured[{ext.lstrip('.')}]\"  —  "
                f"Original error: {exc}"
            ) from exc

    # ------------------------------------------------------------------ #
    #  Format dispatchers                                                  #
    # ------------------------------------------------------------------ #

    async def _dispatch(self, content: bytes, ext: str, filename: str) -> str:
        if ext in {".txt", ".md"}:
            return self._parse_text(content)

        if ext == ".pdf":
            return self._parse_pdf(content)

        if ext == ".docx":
            return self._parse_docx(content)

        if ext == ".pptx":
            return self._parse_pptx(content)

        # HTML and any remaining formats → unstructured
        return self._parse_unstructured(content, ext)

    # ------------------------------------------------------------------ #
    #  Format-specific helpers                                             #
    # ------------------------------------------------------------------ #

    def _parse_text(self, content: bytes) -> str:
        try:
            return content.decode("utf-8")
        except UnicodeDecodeError:
            return content.decode("latin-1", errors="ignore")

    def _parse_pdf(self, content: bytes) -> str:
        import io
        import pypdf

        reader = pypdf.PdfReader(io.BytesIO(content))
        pages = [p.extract_text() for p in reader.pages if p.extract_text()]
        if not pages:
            logger.warning("pypdf extracted no text; trying unstructured fallback.")
            return self._parse_unstructured(content, ".pdf")
        return "\n\n".join(pages)

    def _parse_docx(self, content: bytes) -> str:
        """Fast path using python-docx (bundled with unstructured[docx])."""
        import io
        try:
            from docx import Document  # python-docx
            doc = Document(io.BytesIO(content))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n\n".join(paragraphs)
        except ImportError:
            logger.warning("python-docx not available; falling back to unstructured.")
            return self._parse_unstructured(content, ".docx")

    def _parse_pptx(self, content: bytes) -> str:
        """Fast path using python-pptx (bundled with unstructured[pptx])."""
        import io
        try:
            from pptx import Presentation  # python-pptx
            prs = Presentation(io.BytesIO(content))
            slides: list[str] = []
            for slide in prs.slides:
                texts = [
                    shape.text
                    for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text.strip()
                ]
                if texts:
                    slides.append("\n".join(texts))
            return "\n\n".join(slides)
        except ImportError:
            logger.warning("python-pptx not available; falling back to unstructured.")
            return self._parse_unstructured(content, ".pptx")

    def _parse_unstructured(self, content: bytes, ext: str) -> str:
        """Generic fallback via unstructured.partition.auto."""
        from unstructured.partition.auto import partition

        with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as f:
            f.write(content)
            tmp_path = f.name

        elements = partition(tmp_path)
        return "\n\n".join(str(e) for e in elements)